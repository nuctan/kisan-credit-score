import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 Widescreen format
prs.slide_height = Inches(7.5)

# Color Palette (Warm Saffron & Forest Green Agrarian Theme)
SAFFRON = RGBColor(232, 99, 10)     # #E8630A
FOREST_GREEN = RGBColor(45, 106, 79) # #2D6A4F
CREAM_BG = RGBColor(255, 248, 240)   # #FFF8F0
DARK_BROWN = RGBColor(61, 44, 30)    # #3D2C1E
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
ACCENT_GOLD = RGBColor(212, 160, 23) # #D4A017

def add_header(slide, title_text, category_text="KisanAI Project Presentation"):
    # Header shape background banner
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = FOREST_GREEN
    header_box.line.color.rgb = FOREST_GREEN

    tf = header_box.text_frame
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.15)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD

    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)

# Background
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = FOREST_GREEN
bg1.line.color.rgb = FOREST_GREEN

# Content Box
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
tf1 = title_box.text_frame
tf1.word_wrap = True

p_main = tf1.paragraphs[0]
p_main.text = "🌾 KisanAI"
p_main.font.size = Pt(48)
p_main.font.bold = True
p_main.font.color.rgb = ACCENT_GOLD

p_sub = tf1.add_paragraph()
p_sub.text = "Autonomous Agricultural Credit Assessment & Risk Scoring Platform"
p_sub.font.size = Pt(26)
p_sub.font.bold = True
p_sub.font.color.rgb = WHITE

p_desc = tf1.add_paragraph()
p_desc.text = "Ingesting Sentinel-2 Satellite Telemetry, Econometric Mandi Price Forecasting & Multi-Year Succession Planning"
p_desc.font.size = Pt(16)
p_desc.font.color.rgb = CREAM_BG

p_team = tf1.add_paragraph()
p_team.text = "\nPresented By: Tanishq Kanthed (Nuctan), Akshat Srivastava, Radhika Yadav"
p_team.font.size = Pt(16)
p_team.font.bold = True
p_team.font.color.rgb = ACCENT_GOLD

# Slide 2: Problem Statement
slide2 = prs.slides.add_slide(slide_layout)
add_header(slide2, "2. Problem Statement & Agrarian Challenges")
problems = [
    ("🐢 Slow & Costly Inspections", "Traditional bank field audits cost ₹5,000–₹10,000 per farm and take 3–4 weeks."),
    ("👁️ High Subjectivity & Bias", "Physical visual evaluations vary drastically depending on individual bank officers."),
    ("💸 Debt Traps & Real-Estate Bias", "Banks issue loans based on land real-estate value rather than biological crop yield."),
    ("📉 Static Mandi Price Flaw", "Traditional underwriting uses static annual averages, ignoring harvest-month price crashes.")
]
for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 5.9)
    y = Inches(1.5 + (i // 2) * 2.7)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SAFFRON
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = SAFFRON
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_BROWN

# Slide 3: Objectives
slide3 = prs.slides.add_slide(slide_layout)
add_header(slide3, "3. System Objectives & Goals")
objs = [
    "1. Interactive GIS Mapping: Leaflet + Esri satellite polygon field area calculation in Hectares.",
    "2. Real-Time Remote Sensing: Sentinel-2 L2A optical band processing (B08 NIR, B04 Red) for NDVI.",
    "3. Scikit-Learn Mandi Price ML: Ridge Regression forecasting harvest-month modal prices.",
    "4. Weighted Composite Risk: Combining NDVI (45%), IMD Weather (35%), and Soil N-P-K (20%).",
    "5. 60% DSCR Safe Credit Cap: Protecting farmers against over-indebtedness & debt traps.",
    "6. Multilingual Voice AI: Meta LLaMA 3.3 70B RAG Chatbot supporting 5 Indian languages."
]
for i, text in enumerate(objs):
    y = Inches(1.4 + i * 0.9)
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(0.75))
    box.fill.solid()
    box.fill.fore_color.rgb = CREAM_BG
    box.line.color.rgb = FOREST_GREEN
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BROWN

# Slide 4: System Architecture & Tech Stack
slide4 = prs.slides.add_slide(slide_layout)
add_header(slide4, "4. System Methodology & Tech Stack")
stacks = [
    ("🎨 Frontend Tier", "React 18, Vite 6, TailwindCSS v4, Leaflet.js 1.9.4, Web Speech API"),
    ("⚡ Backend Tier", "Pure Python FastAPI (Uvicorn ASGI), PyMongo with in-memory dict fallback"),
    ("🛰️ Satellite Telemetry", "Sentinel Hub API v3 (Sentinel-2 L2A), 10m resolution, Cloud Masking SCL"),
    ("🤖 AI & ML Engine", "Meta LLaMA 3.3 70B (Groq LPU), Scikit-Learn Ridge Regression, RAG Store")
]
for i, (title, desc) in enumerate(stacks):
    x = Inches(0.8 + (i % 2) * 5.9)
    y = Inches(1.6 + (i // 2) * 2.6)
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = FOREST_GREEN
    card.line.color.rgb = FOREST_GREEN
    tf = card.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(15)
    p1.font.color.rgb = WHITE

# Slide 5: Remote Sensing & GIS
slide5 = prs.slides.add_slide(slide_layout)
add_header(slide5, "5. GIS Mapping & Satellite Remote Sensing")
gis_points = [
    ("🗺️ Leaflet + Esri Imagery", "100% free open-source mapping engine without expensive Google Maps API dependencies."),
    ("📐 Spherical Excess Area Formula", "Calculates exact geodesic field area: Area = (R^2/4) * Σ (Δλ) * (2 + sin φ1 + sin φ2)"),
    ("🛰️ Sentinel-2 L2A Optical Bands", "Ingests 10m spatial resolution Band 8 (NIR) & Band 4 (Red) with cloud masking SCL."),
    ("🌱 NDVI Vegetation Formula", "NDVI = (Band 8 NIR - Band 4 Red) / (Band 8 NIR + Band 4 Red) [0.65+ = Healthy]")
]
for i, (title, desc) in enumerate(gis_points):
    y = Inches(1.4 + i * 1.35)
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = SAFFRON
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(17)
    p0.font.bold = True
    p0.font.color.rgb = SAFFRON
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_BROWN

# Slide 6: Credit Scoring Engine Math
slide6 = prs.slides.add_slide(slide_layout)
add_header(slide6, "6. Credit Scoring Engine — 4-Step Mathematical Model")
steps = [
    ("Step 1: Base Revenue", "Base Rev = Area (Ha) × Hist Yield (T/Ha) × 10 × ML Predicted Price (₹/Q)"),
    ("Step 2: Risk Adjustment", "Adjusted Rev = Base Rev × [(NDVI × 0.45) + (Weather × 0.35) + (Soil × 0.20)]"),
    ("Step 3: Multi-Year Rotation", "Calculates succession crop cycles (Wheat → Mung Bean → Rice) across loan tenure"),
    ("Step 4: 60% DSCR Safe Cap", "Safe Loan Cap = Total Combined Tenure Income × 60% (Corporate DSCR Rule)")
]
for i, (title, desc) in enumerate(steps):
    y = Inches(1.4 + i * 1.35)
    box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = CREAM_BG
    box.line.color.rgb = FOREST_GREEN
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(17)
    p0.font.bold = True
    p0.font.color.rgb = FOREST_GREEN
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_BROWN

# Slide 7: Mandi Price ML Model
slide7 = prs.slides.add_slide(slide_layout)
add_header(slide7, "7. Econometric Mandi Price Forecasting (Scikit-Learn ML)")
mandi_cards = [
    ("📈 Scikit-Learn Ridge Regression", "Trains on historical Agmarknet Mandi dataset to model long-term price trajectory."),
    ("🗓️ Harvest-Month Modulus Math", "t_harvest = (t_sow + duration) mod 12 — Evaluates prices exactly when crop is sold."),
    ("📊 Seasonal Indexing Multiplier", "Applies monthly demand index (e.g., March Wheat Harvest Index = 1.05x multiplier)."),
    ("💡 Real Output Example", "Sowing Wheat in Nov (Harvest in March): Base ₹2,949 → ML Trend ₹3,546 × 1.05 = ₹3,723.73 / Quintal.")
]
for i, (title, desc) in enumerate(mandi_cards):
    x = Inches(0.8 + (i % 2) * 5.9)
    y = Inches(1.5 + (i // 2) * 2.7)
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SAFFRON
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(17)
    p0.font.bold = True
    p0.font.color.rgb = SAFFRON
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_BROWN

# Slide 8: AI Chatbot & RAG
slide8 = prs.slides.add_slide(slide_layout)
add_header(slide8, "8. Multilingual Voice AI Chatbot & RAG Engine")
ai_points = [
    ("🤖 Meta LLaMA 3.3 70B on Groq LPU", "Executes sub-300ms ultra-fast inference (>300 tokens/sec) on Groq hardware."),
    ("📜 RAG Government Schemes Store", "Keyword search over 7 verified schemes (PM-KISAN, KCC, PMFBY, PM-KUSUM, etc.)."),
    ("🛡️ Zero Hallucination Prompting", "System prompt injects farm profile, ML credit cap, and scheme facts to ground LLM."),
    ("🎤 Web Speech Voice Integration", "Browser-native voice input in 5 languages (Hindi, English, Marathi, Gujarati, Tamil).")
]
for i, (title, desc) in enumerate(ai_points):
    y = Inches(1.4 + i * 1.35)
    box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = FOREST_GREEN
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(17)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(14)
    p1.font.color.rgb = WHITE

# Slide 9: Future Scope
slide9 = prs.slides.add_slide(slide_layout)
add_header(slide9, "9. Future Scope & Enhancements")
futures = [
    ("🌐 Pan-India Dataset Scaling", "Expanding crop yield datasets and MANDI historical records to all 28 Indian states."),
    ("📡 Sentinel-1 SAR Radar Integration", "Adding Synthetic Aperture Radar (SAR) to penetrate monsoon cloud cover."),
    ("🔗 Blockchain Credit Passport NFTs", "Issuing immutable, satellite-verified credit certificates for bank API integration.")
]
for i, (title, desc) in enumerate(futures):
    y = Inches(1.5 + i * 1.7)
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = CREAM_BG
    card.line.color.rgb = SAFFRON
    card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = SAFFRON
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(15)
    p1.font.color.rgb = DARK_BROWN

# Slide 10: Conclusion
slide10 = prs.slides.add_slide(slide_layout)
bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg10.fill.solid()
bg10.fill.fore_color.rgb = FOREST_GREEN

c_box = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
tf10 = c_box.text_frame
tf10.word_wrap = True

p_c1 = tf10.paragraphs[0]
p_c1.text = "10. Conclusion"
p_c1.font.size = Pt(36)
p_c1.font.bold = True
p_c1.font.color.rgb = ACCENT_GOLD

p_c2 = tf10.add_paragraph()
p_c2.text = "• KisanAI automates agricultural credit risk evaluation in under 30 seconds, replacing costly manual field audits.\n• Combines Sentinel-2 remote sensing, Scikit-Learn price forecasting, and 60% DSCR safe credit caps to protect both banks and farmers.\n\nThank You! Questions & Discussion."
p_c2.font.size = Pt(20)
p_c2.font.color.rgb = WHITE

# Save Presentation
output_path = "/home/nuctan/Desktop/kisaanai/docs/KISAN_AI_PRESENTATION.pptx"
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
